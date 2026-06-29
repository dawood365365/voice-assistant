import os
from pptx import Presentation
from docx import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings

# ── STEP 1: Extract text from all PPTX files ─────────────────────────────────
# Goes through every .pptx file in the slides/ folder
# Extracts all text from every slide in every file
def load_pptx_files(folder_path: str) -> list:
    documents = []

    for filename in os.listdir(folder_path):
        if filename.endswith(".pptx"):
            filepath = os.path.join(folder_path, filename)
            print(f"📂 Loading: {filename}")

            prs = Presentation(filepath)

            # Go through every slide in this file
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []

                # Go through every shape (text box) on the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                # Combine all text from this slide into one string
                if slide_text:
                    full_text = "\n".join(slide_text)
                    documents.append({
                        "text": full_text,
                        "source": filename,
                        "slide": slide_num + 1
                    })

    print(f"✅ Extracted text from {len(documents)} slides total")
    return documents

# ── LOAD DOCX FILES ────────────────────────────────────────────────────────────
# Goes through every .docx file in the slides/ folder
# Extracts text from paragraphs AND tables (since our Jazz doc has tables)
def load_docx_files(folder_path: str) -> list:
    documents = []

    for filename in os.listdir(folder_path):
        if filename.endswith(".docx"):
            filepath = os.path.join(folder_path, filename)
            print(f"📂 Loading: {filename}")

            doc = Document(filepath)
            current_section = "General"
            section_text = []
            section_num = 1

            # Go through paragraphs — headings mark new sections
            for para in doc.paragraphs:
                if not para.text.strip():
                    continue

                # Treat Heading 1 / Heading 2 styles as section breaks
                if para.style.name.startswith("Heading"):
                    # Save the previous section before starting a new one
                    if section_text:
                        documents.append({
                            "text": "\n".join(section_text),
                            "source": filename,
                            "slide": section_num
                        })
                        section_num += 1
                        section_text = []
                    current_section = para.text.strip()
                    section_text.append(current_section)
                else:
                    section_text.append(para.text.strip())

            # Save the last section
            if section_text:
                documents.append({
                    "text": "\n".join(section_text),
                    "source": filename,
                    "slide": section_num
                })

            # Extract tables separately (e.g. package price tables)
            for t_idx, table in enumerate(doc.tables):
                table_text = []
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    table_text.append(row_text)
                if table_text:
                    documents.append({
                        "text": "\n".join(table_text),
                        "source": filename,
                        "slide": f"table-{t_idx + 1}"
                    })

    print(f"✅ Extracted text from {len(documents)} sections/tables total")
    return documents


# ── STEP 2: Split text into chunks ───────────────────────────────────────────
# Each slide's text gets split into smaller chunks
# chunk_size=500 means each chunk is max 500 characters
# chunk_overlap=50 means chunks share 50 characters with the next chunk
# Overlap helps so important info isn't cut off at a boundary
def split_documents(documents: list) -> list:
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )

    chunks = []
    for doc in documents:
        split_texts = splitter.split_text(doc["text"])
        for chunk in split_texts:
            chunks.append({
                "text": chunk,
                "source": doc["source"],
                "slide": doc["slide"]
            })

    print(f"✅ Split into {len(chunks)} chunks")
    return chunks


# ── STEP 3: Create embeddings and store in ChromaDB ──────────────────────────
# HuggingFace embeddings convert each chunk of text into a vector (list of numbers)
# These vectors capture the meaning of the text
# ChromaDB stores these vectors locally in a folder called chroma_db/
def build_vector_store(chunks: list):
    print("⏳ Loading embedding model (first time may take a minute)...")

    # This downloads a small free model from HuggingFace
    # all-MiniLM-L6-v2 is fast and works well for semantic search
    embeddings = HuggingFaceEmbeddings(
        model_name="all-MiniLM-L6-v2"
    )

    # Extract just the text for ChromaDB
    texts = [chunk["text"] for chunk in chunks]

    # Extract metadata (source file and slide number)
    metadatas = [
        {"source": chunk["source"], "slide": str(chunk["slide"])}
        for chunk in chunks
    ]

    print("⏳ Building vector store...")

    # Create ChromaDB vector store
    # persist_directory saves it to disk so we don't rebuild every time
    vectorstore = Chroma.from_texts(
        texts=texts,
        embedding=embeddings,
        metadatas=metadatas,
        persist_directory="chroma_db"
    )

    print("✅ Vector store saved to chroma_db/")
    return vectorstore


if __name__ == "__main__":
    print("🚀 Starting ingestion pipeline...")

    # Load both PPTX and DOCX files from slides/ folder
    pptx_documents = load_pptx_files("slides")
    docx_documents = load_docx_files("slides")

    documents = pptx_documents + docx_documents

    if not documents:
        print("❌ No slides or documents found! Make sure your files are in the slides/ folder")
        exit()

    print(f"📊 Total source documents found: {len(documents)} (PPTX: {len(pptx_documents)}, DOCX: {len(docx_documents)})")

    # Split into chunks
    chunks = split_documents(documents)

    # Build and save vector store
    build_vector_store(chunks)

    print("\n🎉 Done! Knowledge base is ready.")
    print("You can now run agent.py and ask questions from your slides/documents.")